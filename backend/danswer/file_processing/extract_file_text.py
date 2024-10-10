import io
import json
import os
import re
import zipfile
from collections.abc import Callable
from collections.abc import Iterator
from email.parser import Parser as EmailParser
from pathlib import Path
from typing import Any
from typing import Dict
from typing import IO

import chardet
import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

from danswer.configs.constants import DANSWER_METADATA_FILENAME
from danswer.file_processing.html_utils import parse_html_page_basic
from danswer.file_processing.unstructured import get_unstructured_api_key
from danswer.file_processing.unstructured import unstructured_to_text
from danswer.utils.logger import setup_logger

logger = setup_logger()


TEXT_SECTION_SEPARATOR = "\n\n"


PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".yml",
    ".yaml",
]


VALID_FILE_EXTENSIONS = PLAIN_TEXT_FILE_EXTENSIONS + [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
]


def is_text_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in PLAIN_TEXT_FILE_EXTENSIONS)


def get_file_ext(file_path_or_name: str | Path) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    return extension


def check_file_ext_is_valid(ext: str) -> bool:
    return ext in VALID_FILE_EXTENSIONS


def is_text_file(file: IO[bytes]) -> bool:
    """
    checks if the first 1024 bytes only contain printable or whitespace characters
    if it does, then we say its a plaintext file
    """
    raw_data = file.read(1024)
    text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7F})
    return all(c in text_chars for c in raw_data)


def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    file.seek(0)
    return encoding


def is_macos_resource_fork_file(file_name: str) -> bool:
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )


# To include additional metadata in the search index, add a .danswer_metadata.json file
# to the zip file. This file should contain a list of objects with the following format:
# [{ "filename": "file1.txt", "link": "https://example.com/file1.txt" }]
def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any], dict[str, Any]]]:
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        zip_metadata = {}
        try:
            metadata_file_info = zip_file.getinfo(DANSWER_METADATA_FILENAME)
            with zip_file.open(metadata_file_info, "r") as metadata_file:
                try:
                    zip_metadata = json.load(metadata_file)
                    if isinstance(zip_metadata, list):
                        # convert list of dicts to dict of dicts
                        zip_metadata = {d["filename"]: d for d in zip_metadata}
                except json.JSONDecodeError:
                    logger.warn(f"Unable to load {DANSWER_METADATA_FILENAME}")
        except KeyError:
            logger.info(f"No {DANSWER_METADATA_FILENAME} file")

        for file_info in zip_file.infolist():
            with zip_file.open(file_info.filename, "r") as file:
                if ignore_dirs and file_info.is_dir():
                    continue

                if (
                    ignore_macos_resource_fork_files
                    and is_macos_resource_fork_file(file_info.filename)
                ) or file_info.filename == DANSWER_METADATA_FILENAME:
                    continue
                yield file_info, file, zip_metadata.get(file_info.filename, {})


def _extract_danswer_metadata(line: str) -> dict | None:
    html_comment_pattern = r"<!--\s*DANSWER_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#DANSWER_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None


def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_danswer_metadata: bool = True,
) -> tuple[str, dict]:
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        if ind == 0:
            metadata_or_none = (
                None if ignore_danswer_metadata else _extract_danswer_metadata(line)
            )
            if metadata_or_none is not None:
                metadata = metadata_or_none
            else:
                file_content_raw += line
        else:
            file_content_raw += line

    return file_content_raw, metadata


def pdf_to_text(file: IO[Any], pdf_pass: str | None = None) -> str:
    """Extract text from a PDF file."""
    # Return only the extracted text from read_pdf_file
    text, _ = read_pdf_file(file, pdf_pass)
    return text


def read_pdf_file(
    file: IO[Any],
    pdf_pass: str | None = None,
) -> tuple[str, dict]:
    metadata: Dict[str, Any] = {}
    try:
        pdf_reader = PdfReader(file)

        # If marked as encrypted and a password is provided, try to decrypt
        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            if pdf_pass is not None:
                try:
                    decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
                except Exception:
                    logger.error("Unable to decrypt pdf")

            if not decrypt_success:
                # By user request, keep files that are unreadable just so they
                # can be discoverable by title.
                return "", metadata
        elif pdf_reader.is_encrypted:
            logger.warning("No Password available to decrypt pdf, returning empty")
            return "", metadata

        # Extract metadata from the PDF, removing leading '/' from keys if present
        # This standardizes the metadata keys for consistency
        metadata = {}
        if pdf_reader.metadata is not None:
            for key, value in pdf_reader.metadata.items():
                clean_key = key.lstrip("/")
                if isinstance(value, str) and value.strip():
                    metadata[clean_key] = value

                elif isinstance(value, list) and all(
                    isinstance(item, str) for item in value
                ):
                    metadata[clean_key] = ", ".join(value)

        return (
            TEXT_SECTION_SEPARATOR.join(
                page.extract_text() for page in pdf_reader.pages
            ),
            metadata,
        )
    except PdfStreamError:
        logger.exception("PDF file is not a valid PDF")
    except Exception:
        logger.exception("Failed to read PDF")

    # File is still discoverable by title
    # but the contents are not included as they cannot be parsed
    return "", metadata


def docx_to_text(file: IO[Any]) -> str:
    def is_simple_table(table: docx.table.Table) -> bool:
        for row in table.rows:
            # No omitted cells
            if row.grid_cols_before > 0 or row.grid_cols_after > 0:
                return False

            # No nested tables
            if any(cell.tables for cell in row.cells):
                return False

        return True

    def extract_cell_text(cell: docx.table._Cell) -> str:
        cell_paragraphs = [para.text.strip() for para in cell.paragraphs]
        return " ".join(p for p in cell_paragraphs if p) or "N/A"

    paragraphs = []
    doc = docx.Document(file)
    for item in doc.iter_inner_content():
        if isinstance(item, docx.text.paragraph.Paragraph):
            paragraphs.append(item.text)

        elif isinstance(item, docx.table.Table):
            if not item.rows or not is_simple_table(item):
                continue

            # Every row is a new line, joined with a single newline
            table_content = "\n".join(
                [
                    ",\t".join(extract_cell_text(cell) for cell in row.cells)
                    for row in item.rows
                ]
            )
            paragraphs.append(table_content)

    # Docx already has good spacing between paragraphs
    return "\n".join(paragraphs)


def pptx_to_text(file: IO[Any]) -> str:
    presentation = pptx.Presentation(file)
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        extracted_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
        text_content.append(extracted_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def xlsx_to_text(file: IO[Any]) -> str:
    workbook = openpyxl.load_workbook(file)
    text_content = []
    for sheet in workbook.worksheets:
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(min_row=1, values_only=True)
        )
        text_content.append(sheet_string)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def eml_to_text(file: IO[Any]) -> str:
    text_file = io.TextIOWrapper(file, encoding=detect_encoding(file))
    parser = EmailParser()
    message = parser.parse(text_file)
    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            text_content.append(part.get_payload())
    return TEXT_SECTION_SEPARATOR.join(text_content)


def epub_to_text(file: IO[Any]) -> str:
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)


def file_io_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    file_content_raw, _ = read_text_file(file, encoding=encoding)
    return file_content_raw


def extract_file_text(
    file: IO[Any],
    file_name: str,
    break_on_unprocessable: bool = True,
    extension: str | None = None,
) -> str:
    extension_to_function: dict[str, Callable[[IO[Any]], str]] = {
        ".pdf": pdf_to_text,
        ".docx": docx_to_text,
        ".pptx": pptx_to_text,
        ".xlsx": xlsx_to_text,
        ".eml": eml_to_text,
        ".epub": epub_to_text,
        ".html": parse_html_page_basic,
    }

    try:
        if get_unstructured_api_key():
            return unstructured_to_text(file, file_name)

        if file_name or extension:
            if extension is not None:
                final_extension = extension
            elif file_name is not None:
                final_extension = get_file_ext(file_name)

            if check_file_ext_is_valid(final_extension):
                return extension_to_function.get(final_extension, file_io_to_text)(file)

        # Either the file somehow has no name or the extension is not one that we recognize
        if is_text_file(file):
            return file_io_to_text(file)

        raise ValueError("Unknown file extension and unknown text encoding")

    except Exception as e:
        if break_on_unprocessable:
            raise RuntimeError(
                f"Failed to process file {file_name or 'Unknown'}: {str(e)}"
            ) from e
        logger.warning(f"Failed to process file {file_name or 'Unknown'}: {str(e)}")
        return ""
