import base64
import io
import os
import tempfile
from datetime import datetime
from datetime import timezone
from typing import Any
from typing import List
from typing import Optional

import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from bs4 import BeautifulSoup
from O365 import Account  # type: ignore
from O365.mailbox import MailBox  # type: ignore
from O365.message import Message  # type: ignore
from O365.message import MessageAttachment  # type: ignore

from danswer.configs.app_configs import INDEX_BATCH_SIZE
from danswer.configs.constants import DocumentSource
from danswer.connectors.cross_connector_utils.file_utils import read_pdf_file
from danswer.connectors.interfaces import GenerateDocumentsOutput
from danswer.connectors.interfaces import LoadConnector
from danswer.connectors.interfaces import PollConnector
from danswer.connectors.interfaces import SecondsSinceUnixEpoch
from danswer.connectors.models import BasicExpertInfo
from danswer.connectors.models import ConnectorMissingCredentialError
from danswer.connectors.models import Document
from danswer.connectors.models import Section
from danswer.utils.logger import setup_logger

UNSUPPORTED_FILE_TYPE_CONTENT = ""  # idea copied from the google drive side of thingsy
# Supported file types
SUPPORTED_FILE_TYPES = [".docx", ".xlsx", ".pptx", ".txt", ".pdf"]

logger = setup_logger()


def get_text_from_xlsx_attachment(attachment_object: MessageAttachment) -> str:
    decoded_content = base64.b64decode(attachment_object.content)
    file_content = decoded_content
    logger.info(f"Proccessing XLXS Attachment: {attachment_object.name}")
    excel_file = io.BytesIO(file_content)
    workbook = openpyxl.load_workbook(excel_file, read_only=True)

    full_text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(min_row=1, values_only=True):
            # Check if the row is not completely empty
            if any(cell is not None for cell in row):
                row_string = ",".join(
                    map(lambda x: str(x) if x is not None else "", row)
                )
                # Speed up indexing and AI can handle not having all the empty cells
                row_string = row_string.replace(",,", "")
                full_text.append(row_string)

    return "\n".join(full_text)


def get_text_from_docx_attachment(attachment_object: MessageAttachment) -> str:
    decoded_content = base64.b64decode(attachment_object.content)
    file_content = decoded_content
    full_text = []
    logger.info(f"Proccessing DOCX Attachment: {attachment_object.name}")
    with tempfile.TemporaryDirectory() as local_path:
        with open(os.path.join(local_path, attachment_object.name), "wb") as local_file:
            local_file.write(file_content)
            doc = docx.Document(local_file.name)
            for para in doc.paragraphs:
                full_text.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text.append(cell.text)

    return "\n".join(full_text)


def get_text_from_pdf_attachment(attachment_object: MessageAttachment) -> str:
    decoded_content = base64.b64decode(attachment_object.content)
    file_content = decoded_content
    logger.info(f"Proccessing PDF Attachment: {attachment_object.name}")
    file_text = read_pdf_file(
        file=io.BytesIO(file_content), file_name=attachment_object.name
    )
    return file_text


def get_text_from_txt_attachment(attachment_object: MessageAttachment) -> str:
    decoded_content = base64.b64decode(attachment_object.content)
    logger.info(f"Proccessing TXT Attachment: {attachment_object.name}")
    text_string = decoded_content.decode("utf-8")
    return text_string


def get_text_from_pptx_attachment(attachment_object: MessageAttachment) -> str:
    decoded_content = base64.b64decode(attachment_object.content)
    file_content = decoded_content
    logger.info(f"Proccessing PPTX Attachment: {attachment_object.name}")
    pptx_stream = io.BytesIO(file_content)

    presentation = pptx.Presentation(pptx_stream)
    extracted_text = ""
    for slide_number, slide in enumerate(presentation.slides, start=1):
        extracted_text += f"\nSlide {slide_number}:\n"

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"

    return extracted_text


class ExchangeConnector(LoadConnector, PollConnector):
    def __init__(
        self,
        batch_size: int = INDEX_BATCH_SIZE,
        exchange_max_poll_size: str = "100",
        exchange_categories: list[str] | None = [],
        exchange_folders: list[str] | None = [],
    ) -> None:
        self.batch_size = batch_size
        self.account: Account | None = None
        self.exchange_categories = exchange_categories
        self.exchange_max_poll_size = int(exchange_max_poll_size)
        self.exchange_folders = exchange_folders

    def load_credentials(self, credentials: dict[str, Any]) -> dict[str, Any] | None:
        aad_app_id = credentials["aad_app_id"]
        aad_app_secret = credentials["aad_app_secret"]
        aad_tenant_id = credentials["aad_tenant_id"]
        aad_user_id = credentials["aad_user_id"]

        add_credentials = (aad_app_id, aad_app_secret)
        account = Account(
            add_credentials,
            auth_flow_type="credentials",
            tenant_id=aad_tenant_id,
            main_resource=aad_user_id,
        )
        if not account.is_authenticated:
            account.authenticate()

        self.account = account
        return None

    def _prune_email_list_by_time(
        self, email_list: list, start: datetime, end: datetime
    ) -> list:
        return [
            email
            for email in email_list
            if email.created_date_time >= start and email.created_date_time <= end
        ]

    def get_all_email_objects(
        self,
        limit: int,
        categories: Optional[List[str]] = None,
        folders: Optional[List[str]] = None,
    ) -> list:
        emails: Message = []
        mailbox = self.account.mailbox()
        # Todo: Add support for pagination

        # Only download at max N x Folder x Category most recently modified emails
        # We use Modified because setting catagory updates the modified date to now.
        # This ensures it will be indexed during next poll
        if limit > 999:
            limit = 999

        if categories is None:
            categories = []

        if folders is None:
            folders = []

        # There is probably a better way to do this
        if len(folders) == 0:
            # No folders and No Categories - Tested
            if len(categories) == 0:
                logger.info("Fetching all emails")
                emails = mailbox.get_messages(
                    limit=limit, order_by="receivedDateTime DESC"
                )
            # No folders but has categories - Tested
            else:
                for category in categories:
                    logger.info(f"Fetching Catagory: {category}")
                    emails_in_category = self.get_email_category(limit, category)
                    emails.extend(emails_in_category)
        else:
            # No categories but has folders - Tested
            if len(categories) == 0:
                for folder_path in folders:
                    folder_id = self.get_folder_id(mailbox, folder_path)
                    if folder_id:
                        logger.info(f"Fetching Folder: {folder_path}")
                        folder = mailbox.get_folder(folder_id=folder_id)
                        emails_in_folder = folder.get_messages(
                            limit=limit, order_by="receivedDateTime DESC"
                        )
                        emails.extend(emails_in_folder)
            # Has both folders and categories - Tested
            else:
                for folder_path in folders:
                    for category in categories:
                        logger.info(
                            f"Fetching Folder: {folder_path} and Catagory: {category}"
                        )
                        emails_in_category = self.get_email_category(
                            limit, category, folder_path
                        )
                        emails.extend(emails_in_category)
        return emails

    def get_email_category(
        self, limit: int, category: str, folder_path: Optional[str] = None
    ) -> list:
        logger.info(f"Fetching Catagory: {category}")
        mailbox = self.account.mailbox()
        if folder_path is None:
            folder = mailbox
        else:
            folder_id = self.get_folder_id(mailbox, folder_path)
            folder = mailbox.get_folder(folder_id=folder_id)

        if folder:
            query = mailbox.new_query().any(
                collection="categories", operation="eq", word=category
            )
            emails = folder.get_messages(
                query=query, limit=limit, order_by="lastmodifiedDateTime DESC"
            )

        return emails

    def get_folder_id(self, mailbox: MailBox, folder_path: str) -> str:
        logger.info(f"Fetching Folder ID for: {folder_path}")
        # Split the path into a list of folders
        folders = folder_path.split("/")
        # Get the inbox folder
        folder = mailbox
        # Iterate through the folders
        for folder_name in folders:
            logger.info(f"Folder: {folder_name}")
            # Get the subfolder
            folder = folder.get_folder(folder_name=folder_name)
        # Return the folder_id
        if folder:
            return folder.folder_id
        else:
            return ""

    def clean_html(self, html: str) -> str:
        # Remove HTML tags
        soup = BeautifulSoup(html, "html.parser")
        # Remove unnecessary tags
        for tag in soup(["style", "script", "img"]):
            tag.decompose()
        clean_text = soup.get_text()
        return clean_text

    def _fetch_from_exchange(
        self, start: datetime | None = None, end: datetime | None = None
    ) -> GenerateDocumentsOutput:
        if self.account is None:
            raise ConnectorMissingCredentialError("Exchange")

        email_object_list = self.get_all_email_objects(
            self.exchange_max_poll_size, self.exchange_categories, self.exchange_folders
        )

        if start is not None and end is not None:
            filtered_email_object_list = []
            for email in email_object_list:
                logger.info(f"Email modified: {email.modified}")
                email_received_date = datetime.utcfromtimestamp(
                    email.modified.timestamp()
                )

                # Check if email_received_date is within the start and end range
                if start <= email_received_date <= end:
                    filtered_email_object_list.append(email)
                    logger.info("Email added to filtered list.")
                else:
                    logger.info("Email not within the date range.")

            email_object_list = filtered_email_object_list

        doc_batch: list[Document] = []
        batch_count = 0
        total_count = 0
        for email in email_object_list:
            doc_batch.append(self.convert_email_to_document(email))
            batch_count += 1
            total_count += 1
            # Reached batch size or end of email list
            if batch_count >= self.batch_size or total_count >= len(email_object_list):
                yield doc_batch
                batch_count = 0
                doc_batch = []

        doc_batch = []
        # Attachments
        batch_count = 0
        total_count = 0
        for email in email_object_list:
            # Save Memory and Only download the attachments if there are any
            if email.has_attachments:
                email.attachments.download_attachments()
                for attachment in email.attachments:
                    if attachment.name.endswith(tuple(SUPPORTED_FILE_TYPES)):
                        doc_batch.append(
                            self.convert_attachment_to_document(attachment, email)
                        )

            batch_count += 1
            total_count += 1
            # Reached batch size or end of email list
            if batch_count >= self.batch_size or total_count >= len(email_object_list):
                yield doc_batch
                batch_count = 0
                doc_batch = []
        # yield doc_batch

    def convert_email_to_document(self, email: Message) -> Document:
        # Get the email
        subject = email.subject

        str_to_address = ""
        if email.to:
            for to_address in email.to:
                str_to_address = str_to_address + to_address.address + ", "
        str_cc_address = ""
        if email.cc:
            for cc_address in email.cc:
                str_cc_address = str_cc_address + cc_address.address + ", "
        str_bcc_address = ""
        if email.bcc:
            for bcc_address in email.bcc:
                str_bcc_address = str_bcc_address + bcc_address.address + ", "
        body = self.clean_html(email.body)
        # body_preview = email.body_preview
        sender = email.sender.address
        date = email.received.strftime("%Y-%m-%d %H:%M:%S")
        categories = ""
        for category in email.categories:
            categories += category + ", "
        importance = email.importance
        attachments = ""
        for attachment in email.attachments:
            attachments += attachment.name + ", "

        # Create a unique name for this email
        symantic_identifier = f"{sender} {subject} {date}"
        # create text from the email details
        email_text = f"Subject: {subject}"
        email_text += f"\n\nSender: {sender}"
        email_text += f"\n\nTo: {str_to_address}"
        if str_cc_address:
            email_text += f"\n\nCC: {str_cc_address}"
        if str_bcc_address:
            email_text += f"\n\nBCC: {str_bcc_address}"
        email_text += f"\n\nDate: {date}"
        if categories:
            email_text += f"\n\nCategories: {categories}"
        email_text += f"\n\nImportance: {importance}"
        if attachments:
            email_text += f"\n\nAttachments: {attachments}"
        email_text += f"\n\nBody: {body}"
        link = f"https://outlook.office.com/owa/?ItemID={email.object_id}&viewmodel=ReadMessageItem&path=&exvsurl=1"

        # Create a document object
        document = Document(
            id=email.object_id,
            sections=[Section(link=link, text=email_text)],
            source=DocumentSource.EXCHANGE,
            semantic_identifier=symantic_identifier,
            doc_updated_at=email.modified.replace(tzinfo=timezone.utc),
            primary_owners=[BasicExpertInfo(email=sender)],
            metadata={},
        )

        # Todo: handle calanders

        return document

    def convert_attachment_to_document(
        self, attachment: MessageAttachment, email: Message
    ) -> Document:
        # Check attachment mime type
        attachment_text = UNSUPPORTED_FILE_TYPE_CONTENT
        # If it is a pdf, convert to text
        if attachment.name.endswith(".pdf"):
            attachment_text = get_text_from_pdf_attachment(attachment)
        elif attachment.name.endswith(".xlsx"):
            attachment_text = get_text_from_xlsx_attachment(attachment)
        elif attachment.name.endswith(".docx"):
            attachment_text = get_text_from_docx_attachment(attachment)
        elif attachment.name.endswith(".txt"):
            attachment_text = get_text_from_txt_attachment(attachment)
        elif attachment.name.endswith(".pptx"):
            attachment_text = get_text_from_pptx_attachment(attachment)

        # Documents with the same name can be sent multiple times so adding recieved date to identity
        date = email.received.strftime("%Y-%m-%d %H:%M:%S")
        sender = email.sender.address
        symantic_identifier = f"Attachment: {attachment.name} - {date}"
        id = attachment.attachment_id
        # Todo, link direct to attachment
        link = f"https://outlook.office.com/owa/?ItemID={email.object_id}&viewmodel=ReadMessageItem&path=&exvsurl=1"
        document = Document(
            id=id,
            sections=[Section(link=link, text=attachment_text)],
            source=DocumentSource.EXCHANGE,
            semantic_identifier=symantic_identifier,
            doc_updated_at=email.modified.replace(tzinfo=timezone.utc),
            primary_owners=[BasicExpertInfo(email=sender)],
            metadata={},
        )

        return document

    def load_from_state(self) -> GenerateDocumentsOutput:
        return self._fetch_from_exchange()

    def poll_source(
        self, start: SecondsSinceUnixEpoch, end: SecondsSinceUnixEpoch
    ) -> GenerateDocumentsOutput:
        start_datetime = datetime.utcfromtimestamp(start)
        end_datetime = datetime.utcfromtimestamp(end)
        return self._fetch_from_exchange(start_datetime, end_datetime)


if __name__ == "__main__":
    connector = ExchangeConnector(
        exchange_max_poll_size=os.environ["EXCHANGE_MAX_POLL_SIZE"],
        exchange_categories=[os.environ["EXCHANGE_CATEGORIES"]],
        exchange_folders=[os.environ["EXCHANGE_FOLDERS"]],
    )

    connector.load_credentials(
        {
            "aad_app_id": os.environ["AAD_APP_ID"],
            "aad_app_secret": os.environ["AAD_APP_SECRET"],
            "aad_tenant_id": os.environ["AAD_TENANT_ID"],
            "aad_user_id": os.environ["AAD_USER_ID"],
        }
    )
    document_batches = connector.load_from_state()
    logger.info(next(document_batches))
